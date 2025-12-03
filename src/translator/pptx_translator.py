from pptx import Presentation
from .base import BaseTranslator


class PptxTranslator(BaseTranslator):
    """Translator for PowerPoint files with basic structure preservation."""

    def translate(self, input_path, output_path, target_language):
        """Translate a PowerPoint file while preserving basic structure."""
        print(f"Translating PowerPoint to {target_language}...")

        # Open the presentation
        prs = Presentation(input_path)

        # Process each slide
        for slide_num, slide in enumerate(prs.slides):
            print(f"Translating slide {slide_num + 1}/{len(prs.slides)}...")

            # Process all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Translate text in text frames
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                original_text = paragraph.text
                                translated = self.translate_text(original_text, target_language)

                                # Replace text while trying to preserve formatting
                                for run in paragraph.runs:
                                    run.text = ""
                                if paragraph.runs:
                                    paragraph.runs[0].text = translated
                                else:
                                    paragraph.text = translated

                # Handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                original_text = cell.text
                                translated = self.translate_text(original_text, target_language)
                                cell.text = translated

        prs.save(output_path)
        print(f"Translation complete! Saved to {output_path}")
